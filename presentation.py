from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Helper function to add a slide with title and content
def add_slide(prs, title_text, content_lines, title_font_size=40, body_font_size=24):
    # Use title and content layout (index 1 in most templates)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(title_font_size)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Set content
    content = slide.placeholders[1]  # Content placeholder
    content.text_frame.clear()  # Clear default text
    for line in content_lines:
        p = content.text_frame.paragraphs[0] if not content.text_frame.paragraphs else content.text_frame.add_paragraph()
        p.text = line["text"]
        p.font.size = Pt(line.get("font_size", body_font_size))
        p.font.bold = line.get("bold", False)
        p.level = line.get("level", 0)  # For bullet point indentation
        p.alignment = PP_ALIGN.LEFT
    
    return slide

# Slide 1: Title Slide
add_slide(
    prs,
    "ICT and You – The Dangers of Social Media",
    [
        {"text": "By: The Director of ICT, Achievers University, Owo", "font_size": 28},
        {"text": "Occasion: 2025/2026 Freshmen Orientation Programme", "font_size": 28},
        {"text": "Duration: ≈10 minutes", "font_size": 28}
    ],
    title_font_size=40,
    body_font_size=28
)

# Slide 2: Welcome to the Digital Campus
add_slide(
    prs,
    "1. Welcome to the Digital Campus",
    [
        {"text": "Welcome to Achievers University, a community driven by innovation, discipline, and excellence.", "font_size": 24},
        {"text": "Technology is integral to our learning culture: LMS, Google Workspace, AU-FAP platforms.", "font_size": 24},
        {"text": "ICT offers opportunities but also exposes us to dangers, especially via social media.", "font_size": 24}
    ]
)

# Slide 3: The Blessing and the Burden
add_slide(
    prs,
    "2. The Blessing and the Burden",
    [
        {"text": "Social media (WhatsApp, X, Instagram, TikTok, etc.) revolutionizes communication and learning.", "font_size": 24},
        {"text": "It’s a double-edged sword: a tool to build or break, depending on usage.", "font_size": 24}
    ]
)

# Slide 4: The Hidden Dangers
add_slide(
    prs,
    "3. The Hidden Dangers",
    [
        {"text": "a) Distraction", "font_size": 24, "bold": True},
        {"text": "Endless scrolling steals study time. Discipline is key to success.", "font_size": 20, "level": 1},
        {"text": "b) Digital Footprints", "font_size": 24, "bold": True},
        {"text": "Posts can resurface and affect future opportunities. The internet never forgets.", "font_size": 20, "level": 1},
        {"text": "c) Cybercrime", "font_size": 24, "bold": True},
        {"text": "Beware of scams, identity theft, and blackmail. Protect personal data.", "font_size": 20, "level": 1},
        {"text": "d) Misinformation", "font_size": 24, "bold": True},
        {"text": "Fake news spreads fast. Verify before sharing.", "font_size": 20, "level": 1},
        {"text": "e) Mental Health", "font_size": 24, "bold": True},
        {"text": "Comparison can harm self-esteem. Follow positive, educational content.", "font_size": 20, "level": 1}
    ]
)

# Slide 5: ICT Responsibility
add_slide(
    prs,
    "4. ICT Responsibility at Achievers University",
    [
        {"text": "ICT is about responsible digital citizenship.", "font_size": 24},
        {"text": "Our policy promotes:", "font_size": 24},
        {"text": "Academic integrity", "font_size": 20, "level": 1},
        {"text": "Cybersecurity awareness", "font_size": 20, "level": 1},
        {"text": "Responsible digital communication", "font_size": 20, "level": 1},
        {"text": "Smart use of educational tools", "font_size": 20, "level": 1}
    ]
)

# Slide 6: Using Social Media Wisely
add_slide(
    prs,
    "5. Using Social Media Wisely",
    [
        {"text": "Five principles for safe online presence:", "font_size": 24},
        {"text": "Pause before you post", "font_size": 20, "level": 1},
        {"text": "Protect your privacy", "font_size": 20, "level": 1},
        {"text": "Promote positivity", "font_size": 20, "level": 1},
        {"text": "Prioritize your purpose", "font_size": 20, "level": 1},
        {"text": "Participate in learning", "font_size": 20, "level": 1}
    ]
)

# Slide 7: Conclusion
add_slide(
    prs,
    "6. Conclusion",
    [
        {"text": "ICT offers immense opportunities but requires responsibility.", "font_size": 24},
        {"text": "Social media can make or mar you. Choose wisely.", "font_size": 24},
        {"text": "Be wise, disciplined, and digitally responsible.", "font_size": 24},
        {"text": "Welcome to Achievers University — where technology meets purpose!", "font_size": 24, "bold": True}
    ]
)

# Save the presentation
prs.save("ICT_and_You.pptx")
print("PowerPoint presentation saved as 'ICT_and_You.pptx'")